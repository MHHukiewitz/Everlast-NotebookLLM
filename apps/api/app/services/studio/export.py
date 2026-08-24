import csv
import html
import io
import json
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.services.pdf import AI_MARK, markdown_to_pdf
from app.services.studio.png import infographic_png, mindmap_png

FORMATS: dict[str, tuple[str, ...]] = {
    "note": ("md", "txt", "pdf", "json"),
    "report": ("md", "pdf", "json"),
    "mindmap": ("png", "mmd", "pdf", "json"),
    "quiz": ("md", "csv", "pdf", "json"),
    "flashcards": ("md", "csv", "pdf", "json"),
    "table": ("csv", "md", "pdf", "json"),
    "slides": ("md", "pptx", "html", "pdf", "txt", "json"),
    "infographic": ("png", "svg", "md", "pdf", "json"),
    "audio": ("mp3", "md", "json"),
    "video": ("mp4", "md", "json"),
}

MEDIA: dict[str, str] = {
    "md": "text/markdown; charset=utf-8",
    "txt": "text/plain; charset=utf-8",
    "pdf": "application/pdf",
    "json": "application/json; charset=utf-8",
    "csv": "text/csv; charset=utf-8",
    "mmd": "text/plain; charset=utf-8",
    "svg": "image/svg+xml",
    "png": "image/png",
    "html": "text/html; charset=utf-8",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "mp3": "audio/mpeg",
    "mp4": "video/mp4",
}


def allowed_formats(artifact_type: str) -> tuple[str, ...]:
    return FORMATS.get(artifact_type, ())


def artifact_markdown(artifact_type: str, title: str, payload: dict[str, Any]) -> str:
    if artifact_type == "note":
        return _note_md(title, payload)
    if artifact_type == "report":
        return _report_md(title, payload)
    if artifact_type == "mindmap":
        return _mindmap_md(title, payload)
    if artifact_type == "quiz":
        return _quiz_md(title, payload)
    if artifact_type == "flashcards":
        return _flashcards_md(title, payload)
    if artifact_type == "table":
        return _table_md(title, payload)
    if artifact_type == "slides":
        return _slides_md(title, payload)
    if artifact_type == "infographic":
        return _infographic_md(title, payload)
    if artifact_type == "audio":
        return _audio_md(title, payload)
    if artifact_type == "video":
        return _video_md(title, payload)
    raise ValueError("Dieses Artefakt hat keinen Text.")


def safe_name(title: str) -> str:
    cleaned = "".join(char if char.isalnum() or char in "-_ " else "_" for char in title)
    return (cleaned.strip() or "studio")[:40]


def _text(value: str) -> bytes:
    return value.encode("utf-8")


def _json_bytes(artifact_type: str, title: str, payload: dict[str, Any]) -> bytes:
    body = {
        "ai_generated": True,
        "notice": AI_MARK,
        "type": artifact_type,
        "title": title,
        "payload": payload,
    }
    return json.dumps(body, ensure_ascii=False, indent=2).encode("utf-8")


def _csv_bytes(rows: list[list[str]]) -> bytes:
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    for row in rows:
        writer.writerow(row)
    return buffer.getvalue().encode("utf-8")


def _note_md(title: str, payload: dict[str, Any]) -> str:
    return f"# {title}\n\n{payload.get('body') or ''}\n\n{AI_MARK}\n"


def _report_md(title: str, payload: dict[str, Any]) -> str:
    return f"# {title}\n\n{payload.get('body_md') or ''}\n\n{AI_MARK}\n"


def _mindmap_md(title: str, payload: dict[str, Any]) -> str:
    mermaid = payload.get("mermaid") or ""
    return f"# {title}\n\n```mermaid\n{mermaid}\n```\n\n{AI_MARK}\n"


def _quiz_md(title: str, payload: dict[str, Any]) -> str:
    parts = [f"# {title}", ""]
    for index, question in enumerate(payload.get("questions") or [], start=1):
        parts.append(f"## {index}. {question.get('question') or ''}")
        for choice in question.get("choices") or []:
            parts.append(f"- {choice}")
        answer_index = question.get("answer_index")
        choices = question.get("choices") or []
        if isinstance(answer_index, int) and 0 <= answer_index < len(choices):
            parts.append(f"Antwort: {choices[answer_index]}")
        if question.get("explanation"):
            parts.append(str(question.get("explanation")))
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _quiz_csv(payload: dict[str, Any]) -> list[list[str]]:
    rows = [["front", "back"]]
    for question in payload.get("questions") or []:
        choices = question.get("choices") or []
        answer_index = question.get("answer_index")
        answer = ""
        if isinstance(answer_index, int) and 0 <= answer_index < len(choices):
            answer = str(choices[answer_index])
        rows.append([str(question.get("question") or ""), answer])
    rows.append(["(KI-generiert)", AI_MARK])
    return rows


def _flashcards_md(title: str, payload: dict[str, Any]) -> str:
    parts = [f"# {title}", ""]
    for card in payload.get("cards") or []:
        parts.append(f"## {card.get('front') or ''}")
        parts.append(str(card.get("back") or ""))
        if card.get("cite"):
            parts.append(str(card.get("cite")))
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _flashcards_csv(payload: dict[str, Any]) -> list[list[str]]:
    rows = [["front", "back"]]
    for card in payload.get("cards") or []:
        back = str(card.get("back") or "")
        if card.get("cite"):
            back = f"{back} {card.get('cite')}"
        rows.append([str(card.get("front") or ""), back])
    rows.append(["(KI-generiert)", AI_MARK])
    return rows


def _table_md(title: str, payload: dict[str, Any]) -> str:
    columns = [str(column) for column in payload.get("columns") or []]
    rows = payload.get("rows") or []
    parts = [f"# {title}", ""]
    if columns:
        parts.append("| " + " | ".join(columns) + " |")
        parts.append("| " + " | ".join("---" for _ in columns) + " |")
        for row in rows:
            cells = [str(cell) for cell in row]
            while len(cells) < len(columns):
                cells.append("")
            parts.append("| " + " | ".join(cells[: len(columns)]) + " |")
    parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _xml_escape(value: str) -> str:
    return (
        value.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _slides_md(title: str, payload: dict[str, Any]) -> str:
    parts = [f"# {title}", ""]
    for index, slide in enumerate(payload.get("slides") or [], start=1):
        parts.append(f"## {index}. {slide.get('heading') or ''}")
        for bullet in slide.get("bullets") or []:
            parts.append(f"- {bullet}")
        if slide.get("notes"):
            parts.append(str(slide.get("notes")))
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _slides_txt(title: str, payload: dict[str, Any]) -> str:
    parts = [title, ""]
    for index, slide in enumerate(payload.get("slides") or [], start=1):
        parts.append(f"{index}. {slide.get('heading') or ''}")
        for bullet in slide.get("bullets") or []:
            parts.append(f"- {bullet}")
        if slide.get("notes"):
            parts.append(f"Notizen: {slide.get('notes')}")
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _slides_html(title: str, payload: dict[str, Any]) -> str:
    slides = payload.get("slides") or []
    sections: list[str] = []
    sections.append(
        "<section class=\"slide title-slide\">"
        f"<h1>{html.escape(title)}</h1>"
        f"<p class=\"mark\">{html.escape(AI_MARK)}</p>"
        "</section>"
    )
    for index, slide in enumerate(slides, start=1):
        bullets = "".join(
            f"<li>{html.escape(str(bullet))}</li>" for bullet in slide.get("bullets") or []
        )
        notes = slide.get("notes")
        notes_html = f"<p class=\"notes\">{html.escape(str(notes))}</p>" if notes else ""
        sections.append(
            "<section class=\"slide\">"
            f"<p class=\"count\">{index} / {len(slides)}</p>"
            f"<h2>{html.escape(str(slide.get('heading') or ''))}</h2>"
            f"<ul>{bullets}</ul>"
            f"{notes_html}"
            f"<p class=\"mark\">{html.escape(AI_MARK)}</p>"
            "</section>"
        )
    return f"""<!doctype html>
<html lang="de">
<head>
  <meta charset="utf-8" />
  <meta name="generator" content="Everlast Notebook" />
  <meta name="X-AI-Generated" content="true" />
  <title>{html.escape(title)}</title>
  <style>
    html, body {{ margin: 0; background: #111; color: #1f1f1f; font-family: Inter, ui-sans-serif, sans-serif; }}
    .deck {{ height: 100vh; overflow: hidden; }}
    .slide {{
      box-sizing: border-box; height: 100vh; padding: 56px 72px 72px;
      background: #fff; display: none; flex-direction: column;
    }}
    .slide.current {{ display: flex; }}
    h1 {{ font-size: 48px; margin: auto 0; }}
    h2 {{ font-size: 36px; margin: 0 0 24px; }}
    ul {{ font-size: 24px; line-height: 1.45; margin: 0; padding-left: 28px; }}
    li {{ margin: 0 0 12px; }}
    .count {{ color: #737373; font-size: 14px; margin: 0 0 12px; }}
    .notes {{ color: #525252; font-size: 16px; margin-top: auto; }}
    .mark {{ color: #737373; font-size: 12px; margin-top: auto; }}
    .title-slide .mark {{ margin-top: 24px; }}
  </style>
</head>
<body>
  <div class="deck">
    {"".join(sections)}
  </div>
  <script>
    const slides = Array.from(document.querySelectorAll(".slide"));
    let index = 0;
    function show(next) {{
      index = Math.max(0, Math.min(slides.length - 1, next));
      slides.forEach((slide, i) => slide.classList.toggle("current", i === index));
    }}
    show(0);
    document.addEventListener("keydown", (event) => {{
      if (event.key === "ArrowRight" || event.key === " ") show(index + 1);
      if (event.key === "ArrowLeft") show(index - 1);
    }});
  </script>
</body>
</html>
"""


def _pptx_box(slide: Any, left: Any, top: Any, width: Any, height: Any, text: str, size: int, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(31, 31, 31)


def _slides_pptx(title: str, payload: dict[str, Any]) -> bytes:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]
    title_slide = deck.slides.add_slide(blank)
    _pptx_box(title_slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.6), title, 40, True)
    _pptx_box(title_slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4), AI_MARK, 12)
    for slide_data in payload.get("slides") or []:
        slide = deck.slides.add_slide(blank)
        heading = str(slide_data.get("heading") or "")
        _pptx_box(slide, Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.1), heading, 32, True)
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.6))
        frame = body.text_frame
        frame.word_wrap = True
        bullets = [str(bullet) for bullet in slide_data.get("bullets") or []]
        if bullets:
            frame.paragraphs[0].text = bullets[0]
            frame.paragraphs[0].font.size = Pt(20)
            for bullet in bullets[1:]:
                paragraph = frame.add_paragraph()
                paragraph.text = bullet
                paragraph.font.size = Pt(20)
                paragraph.space_before = Pt(8)
        notes = str(slide_data.get("notes") or "").strip()
        slide.notes_slide.notes_text_frame.text = f"{notes}\n\n{AI_MARK}".strip()
        _pptx_box(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.35), AI_MARK, 10)
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _audio_md(title: str, payload: dict[str, Any]) -> str:
    parts = [f"# {title}", ""]
    for turn in payload.get("turns") or []:
        text = str(turn.get("text") or "").strip()
        if not text:
            continue
        parts.append(text)
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _video_md(title: str, payload: dict[str, Any]) -> str:
    parts = [f"# {title}", ""]
    for index, scene in enumerate(payload.get("scenes") or [], start=1):
        parts.append(f"## {index}. {scene.get('heading') or ''}")
        for bullet in scene.get("bullets") or []:
            parts.append(f"- {bullet}")
        if scene.get("narration"):
            parts.append(str(scene.get("narration")))
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _media_bytes(payload: dict[str, Any], key: str) -> bytes:
    path = payload.get(key)
    if not path:
        raise ValueError("Die Mediendatei ist noch nicht fertig.")
    file_path = Path(str(path))
    if not file_path.exists():
        raise ValueError("Die Mediendatei ist noch nicht fertig.")
    return file_path.read_bytes()


def _infographic_md(title: str, payload: dict[str, Any]) -> str:
    parts = [f"# {title}", ""]
    for chart in payload.get("charts") or []:
        parts.append(f"## {chart.get('title') or 'Diagramm'}")
        unit = chart.get("unit") or ""
        for point in chart.get("points") or []:
            suffix = f" {unit}" if unit else ""
            parts.append(f"- {point.get('label')}: {point.get('value')}{suffix}")
        parts.append("")
    for item in payload.get("items") or []:
        parts.append(f"## {item.get('label') or ''}")
        if item.get("number"):
            parts.append(f"{item.get('number')}{item.get('suffix') or ''}")
        if item.get("caption"):
            parts.append(str(item.get("caption")))
        if item.get("value") and item.get("value") != item.get("caption"):
            parts.append(str(item.get("value")))
        if item.get("detail") and item.get("detail") != item.get("caption"):
            parts.append(str(item.get("detail")))
        parts.append("")
    parts.append(AI_MARK)
    return "\n".join(parts) + "\n"


def _wrap_words(text: str, width: int) -> list[str]:
    words = str(text or "").split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        trial = f"{current} {word}"
        if len(trial) <= width:
            current = trial
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines[:6]


CHART_COLORS = ("#2563eb", "#7c3aed", "#16a34a", "#d97706", "#db2777", "#0891b2")


def _chart_svg_blocks(charts: list[Any], top: int) -> tuple[str, int]:
    blocks: list[str] = []
    y = top
    for chart in charts:
        points = [point for point in chart.get("points") or [] if isinstance(point, dict)]
        if len(points) < 2:
            continue
        title = str(chart.get("title") or "Diagramm")
        kind = str(chart.get("type") or "bar")
        unit = str(chart.get("unit") or "")
        blocks.append(
            f'<text x="16" y="{y + 18}" font-size="14" font-weight="600" fill="#1f1f1f">{_xml_escape(title)}</text>'
        )
        y += 28
        values = [float(point.get("value") or 0) for point in points]
        peak = max(values) or 1
        if kind == "pie":
            cx, cy, radius = 140, y + 90, 70
            total = sum(values) or 1
            angle = -math.pi / 2
            for index, point in enumerate(points):
                sweep = (values[index] / total) * 2 * math.pi
                end = angle + sweep
                x1 = cx + radius * math.cos(angle)
                y1 = cy + radius * math.sin(angle)
                x2 = cx + radius * math.cos(end)
                y2 = cy + radius * math.sin(end)
                large = 1 if sweep > math.pi else 0
                color = CHART_COLORS[index % len(CHART_COLORS)]
                blocks.append(
                    f'<path d="M {cx:.1f} {cy:.1f} L {x1:.1f} {y1:.1f} A {radius} {radius} 0 {large} 1 {x2:.1f} {y2:.1f} Z" fill="{color}"/>'
                )
                blocks.append(
                    f'<rect x="230" y="{y + 8 + index * 22}" width="10" height="10" fill="{color}"/>'
                    f'<text x="246" y="{y + 18 + index * 22}" font-size="11" fill="#1f1f1f">{_xml_escape(str(point.get("label") or ""))} · {values[index]:g}{(" " + unit) if unit else ""}</text>'
                )
                angle = end
            y += 200
        elif kind == "bar":
            width = max(24, int(560 / len(points)) - 8)
            for index, point in enumerate(points):
                bar_h = max(4, int((values[index] / peak) * 110))
                x = 24 + index * (width + 8)
                color = CHART_COLORS[index % len(CHART_COLORS)]
                blocks.append(
                    f'<rect x="{x}" y="{y + 120 - bar_h}" width="{width}" height="{bar_h}" rx="4" fill="{color}"/>'
                    f'<text x="{x}" y="{y + 138}" font-size="10" fill="#525252">{_xml_escape(str(point.get("label") or "")[:12])}</text>'
                    f'<text x="{x}" y="{y + 16}" font-size="11" fill="#1f1f1f">{values[index]:g}</text>'
                )
            y += 160
        else:
            for index, point in enumerate(points):
                bar_w = max(8, int((values[index] / peak) * 360))
                color = CHART_COLORS[index % len(CHART_COLORS)]
                yy = y + index * 28
                blocks.append(
                    f'<text x="16" y="{yy + 14}" font-size="11" fill="#1f1f1f">{_xml_escape(str(point.get("label") or ""))}</text>'
                    f'<rect x="200" y="{yy}" width="{bar_w}" height="16" rx="4" fill="{color}"/>'
                    f'<text x="{208 + bar_w}" y="{yy + 13}" font-size="11" fill="#1f1f1f">{values[index]:g}{(" " + unit) if unit else ""}</text>'
                )
            y += len(points) * 28 + 16
        y += 12
    return "".join(blocks), y


def infographic_svg(title: str, payload: dict[str, Any]) -> str:
    items = payload.get("items") or []
    fills = ("#eff6ff", "#f5f3ff", "#f0fdf4", "#fff7ed")
    chart_svg, y = _chart_svg_blocks(payload.get("charts") or [], 56)
    cards: list[tuple[int, int, int, dict[str, Any], list[str], list[str], list[str]]] = []
    col_w = 300
    gap = 16
    row_h = 0
    for index, item in enumerate(items):
        caption = str(item.get("caption") or "")
        value = str(item.get("value") or "")
        detail = str(item.get("detail") or "")
        caption_lines = _wrap_words(caption, 34) if caption else []
        value_lines = _wrap_words(value, 28) if value and value != caption else []
        detail_lines = _wrap_words(detail, 34) if detail and detail != caption else []
        number_h = 36 if item.get("number") else 0
        height = 56 + number_h + 16 * len(caption_lines) + 20 * len(value_lines) + 16 * len(detail_lines)
        col = index % 2
        if col == 0:
            row_h = height
        else:
            row_h = max(row_h, height)
        x = 16 + col * (col_w + gap)
        cards.append((x, y, height, item, caption_lines, value_lines, detail_lines))
        if col == 1 or index == len(items) - 1:
            y += row_h + gap
    total_h = max(y + 28, 160)
    blocks: list[str] = []
    for index, (x, top, height, item, caption_lines, value_lines, detail_lines) in enumerate(cards):
        fill = fills[index % len(fills)]
        cursor = top + 22
        number_svg = []
        if item.get("number"):
            cursor += 32
            suffix = str(item.get("suffix") or "")
            suffix_svg = (
                f'<tspan font-size="20" font-weight="500" fill="#525252">{_xml_escape(suffix)}</tspan>'
                if suffix
                else ""
            )
            gap_svg = " " if suffix_svg else ""
            number_svg.append(
                f'<text x="{x + 16}" y="{cursor}" font-size="28" font-weight="600" fill="#1f1f1f">{_xml_escape(str(item.get("number")))}{gap_svg}{suffix_svg}</text>'
            )
        caption_svg = []
        for line in caption_lines:
            cursor += 16
            caption_svg.append(
                f'<text x="{x + 16}" y="{cursor}" font-size="12" fill="#404040">{_xml_escape(line)}</text>'
            )
        value_svg = []
        for line in value_lines:
            cursor += 20
            value_svg.append(
                f'<text x="{x + 16}" y="{cursor}" font-size="15" font-weight="600" fill="#1f1f1f">{_xml_escape(line)}</text>'
            )
        detail_svg = []
        for line in detail_lines:
            cursor += 16
            detail_svg.append(
                f'<text x="{x + 16}" y="{cursor}" font-size="11" fill="#525252">{_xml_escape(line)}</text>'
            )
        blocks.append(
            f'<rect x="{x}" y="{top}" width="{col_w}" height="{height}" rx="10" fill="{fill}" stroke="#e5e5e5"/>'
            f'<text x="{x + 16}" y="{top + 22}" font-size="11" fill="#737373">{_xml_escape(str(index + 1) + " · " + str(item.get("label") or ""))}</text>'
            + "".join(number_svg)
            + "".join(caption_svg)
            + "".join(value_svg)
            + "".join(detail_svg)
        )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 640 {total_h}" width="640" height="{total_h}">'
        f'<rect width="640" height="{total_h}" fill="#ffffff"/>'
        f'<text x="16" y="36" font-size="18" font-weight="600" fill="#1f1f1f">{_xml_escape(title)}</text>'
        f'<text x="16" y="{total_h - 10}" font-size="10" fill="#737373">{_xml_escape(AI_MARK)}</text>'
        + chart_svg
        + "".join(blocks)
        + "</svg>"
    )


def _table_csv(payload: dict[str, Any]) -> list[list[str]]:
    columns = [str(column) for column in payload.get("columns") or []]
    rows = [columns]
    for row in payload.get("rows") or []:
        rows.append([str(cell) for cell in row])
    notice = [AI_MARK]
    notice.extend("" for _ in columns[1:])
    rows.append(notice)
    return rows


def export_artifact(
    artifact_type: str, title: str, payload: dict[str, Any], fmt: str
) -> tuple[bytes, str, str]:
    allowed = allowed_formats(artifact_type)
    if not allowed or fmt not in allowed:
        raise ValueError("Dieses Exportformat ist nicht verfügbar.")
    if fmt == "json":
        data = _json_bytes(artifact_type, title, payload)
    elif artifact_type == "note" and fmt == "md":
        data = _text(_note_md(title, payload))
    elif artifact_type == "note" and fmt == "txt":
        data = _text(f"{title}\n\n{payload.get('body') or ''}\n\n{AI_MARK}\n")
    elif artifact_type == "note" and fmt == "pdf":
        data = markdown_to_pdf(title, _note_md(title, payload))
    elif artifact_type == "report" and fmt == "md":
        data = _text(_report_md(title, payload))
    elif artifact_type == "report" and fmt == "pdf":
        data = markdown_to_pdf(title, _report_md(title, payload))
    elif artifact_type == "mindmap" and fmt == "png":
        data = mindmap_png(title, payload)
    elif artifact_type == "mindmap" and fmt == "mmd":
        data = _text(f"%% {AI_MARK}\n{payload.get('mermaid') or ''}\n")
    elif artifact_type == "mindmap" and fmt == "pdf":
        data = markdown_to_pdf(title, _mindmap_md(title, payload))
    elif artifact_type == "quiz" and fmt == "md":
        data = _text(_quiz_md(title, payload))
    elif artifact_type == "quiz" and fmt == "csv":
        data = _csv_bytes(_quiz_csv(payload))
    elif artifact_type == "quiz" and fmt == "pdf":
        data = markdown_to_pdf(title, _quiz_md(title, payload))
    elif artifact_type == "flashcards" and fmt == "md":
        data = _text(_flashcards_md(title, payload))
    elif artifact_type == "flashcards" and fmt == "csv":
        data = _csv_bytes(_flashcards_csv(payload))
    elif artifact_type == "flashcards" and fmt == "pdf":
        data = markdown_to_pdf(title, _flashcards_md(title, payload))
    elif artifact_type == "table" and fmt == "md":
        data = _text(_table_md(title, payload))
    elif artifact_type == "table" and fmt == "csv":
        data = _csv_bytes(_table_csv(payload))
    elif artifact_type == "table" and fmt == "pdf":
        data = markdown_to_pdf(title, _table_md(title, payload))
    elif artifact_type == "slides" and fmt == "md":
        data = _text(_slides_md(title, payload))
    elif artifact_type == "slides" and fmt == "pptx":
        data = _slides_pptx(title, payload)
    elif artifact_type == "slides" and fmt == "html":
        data = _text(_slides_html(title, payload))
    elif artifact_type == "slides" and fmt == "pdf":
        data = markdown_to_pdf(title, _slides_md(title, payload))
    elif artifact_type == "slides" and fmt == "txt":
        data = _text(_slides_txt(title, payload))
    elif artifact_type == "infographic" and fmt == "png":
        data = infographic_png(title, payload)
    elif artifact_type == "infographic" and fmt == "svg":
        data = _text(infographic_svg(title, payload))
    elif artifact_type == "infographic" and fmt == "md":
        data = _text(_infographic_md(title, payload))
    elif artifact_type == "infographic" and fmt == "pdf":
        data = markdown_to_pdf(title, _infographic_md(title, payload))
    elif artifact_type == "audio" and fmt == "md":
        data = _text(_audio_md(title, payload))
    elif artifact_type == "audio" and fmt == "mp3":
        data = _media_bytes(payload, "audio_path")
    elif artifact_type == "video" and fmt == "md":
        data = _text(_video_md(title, payload))
    elif artifact_type == "video" and fmt == "mp4":
        data = _media_bytes(payload, "video_path")
    else:
        raise ValueError("Dieses Exportformat ist nicht verfügbar.")
    filename = f"{safe_name(title)}.{fmt}"
    return data, MEDIA[fmt], filename
